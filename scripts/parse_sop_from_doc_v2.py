"""SOP Parser v2 — Extracts BOTH the main SOP table AND all embedded OLE objects.

Extends v1 by:
  1. Extracting embedded PowerPoint (.pptx) files → Acceptable/Supplement document lists
  2. Extracting embedded Excel (.xlsx) files → NNS Proceed Checklist decision matrix
  3. Noting embedded PDF objects that require Acrobat for extraction
  4. Merging everything into a single comprehensive sop_rules_parsed_v2.json

Requires: Windows with Microsoft Word, PowerPoint, and Excel installed.
          pip install python-pptx openpyxl

Usage:
    python scripts/parse_sop_from_doc_v2.py
"""

from __future__ import annotations

import json
import os
import re
import time


# ---------------------------------------------------------------------------
# 1. Word COM — extract main SOP table + embedded OLE objects to disk
# ---------------------------------------------------------------------------

def extract_sop_table(doc) -> dict[int, dict[int, str]]:
    """Extract the main SOP exception matrix (Table 2) from an open Word doc."""
    table_data: dict[int, dict[int, str]] = {}
    if doc.Tables.Count >= 2:
        table = doc.Tables(2)
        for r in range(2, table.Rows.Count + 1):  # skip header row
            row: dict[int, str] = {}
            for c in range(1, table.Columns.Count + 1):
                try:
                    cell_text = table.Cell(r, c).Range.Text
                    cell_text = cell_text.replace("\x07", "").replace("\x01", "").strip()
                    row[c] = cell_text
                except Exception:
                    row[c] = ""
            table_data[r] = row
    return table_data


def extract_embedded_objects(doc, out_dir: str) -> dict[str, str]:
    """Activate each embedded OLE object and save it to disk via its native app.

    Returns a dict mapping icon labels to saved file paths.
    """
    import win32com.client

    os.makedirs(out_dir, exist_ok=True)
    extracted: dict[str, str] = {}

    for i in range(1, doc.InlineShapes.Count + 1):
        shape = doc.InlineShapes(i)
        try:
            ole = shape.OLEFormat
        except Exception:
            continue

        label = ole.IconLabel
        prog_id = ole.ProgID
        save_path = os.path.join(out_dir, label)

        print(f"  [{i}] {label} ({prog_id})")

        if "PowerPoint" in prog_id:
            ole.DoVerb(0)
            time.sleep(4)
            try:
                ppt = win32com.client.GetActiveObject("PowerPoint.Application")
                if ppt.Presentations.Count > 0:
                    pres = ppt.ActivePresentation
                    pres.SaveAs(save_path)
                    pres.Close()
                    extracted[label] = save_path
                    print(f"       -> Saved ({os.path.getsize(save_path):,} bytes)")
                try:
                    ppt.Quit()
                except Exception:
                    pass
            except Exception as e:
                print(f"       -> PowerPoint extraction failed: {e}")
            time.sleep(1)

        elif "Excel" in prog_id:
            ole.DoVerb(0)
            time.sleep(4)
            try:
                xl = win32com.client.GetActiveObject("Excel.Application")
                if xl.Workbooks.Count > 0:
                    wb = xl.ActiveWorkbook
                    wb.SaveAs(save_path)
                    wb.Close(False)
                    extracted[label] = save_path
                    print(f"       -> Saved ({os.path.getsize(save_path):,} bytes)")
                try:
                    xl.Quit()
                except Exception:
                    pass
            except Exception as e:
                print(f"       -> Excel extraction failed: {e}")
            time.sleep(1)

        elif "Acrobat" in prog_id or "PDF" in prog_id.upper():
            print("       -> Skipped (PDF requires Acrobat COM; noted in output)")
            extracted[label] = "__PDF_NOT_EXTRACTED__"

    return extracted


# ---------------------------------------------------------------------------
# 2. Parse main SOP table steps (same logic as v1)
# ---------------------------------------------------------------------------

def _parse_steps(raw_text: str) -> list[dict]:
    """Generic step parser: extract Step N / Step XA/XB patterns."""
    steps = re.findall(r"Step\s+(\w+):\s*(.+?)(?=Step|\Z)", raw_text, re.DOTALL)
    result = []
    for step_id, action in steps:
        action = action.strip().rstrip("\r\n")
        entry: dict = {}

        if re.match(r"^\d+$", step_id):
            entry["step"] = int(step_id)
        else:
            entry["branch"] = step_id

        # Conditions
        if "If matched" in action:
            entry["condition"] = "Matched"
        elif "If not matched" in action or "Not matched" in action:
            entry["condition"] = "Not matched"
        elif "If Proceed" in action:
            entry["condition"] = "Proceed"
        elif "If Remediate" in action:
            entry["condition"] = "Remediate first"
        elif "If Not Matched" in action:
            entry["condition"] = "No match"
        elif "If Potential Match" in action:
            entry["condition"] = "Potential match"

        entry["action"] = action

        # Statuses
        if '"Completed"' in action or "\u201cCompleted\u201d" in action:
            entry["status"] = "Completed"
        elif "Supplement" in action:
            entry["status"] = "SupplementRequired"
        elif "Not Required" in action or "Not required" in action:
            entry["status"] = "NotRequired"
        elif "investigation required" in action.lower():
            entry["status"] = "InvestigationRequired"

        result.append(entry)
    return result


def parse_identity(ecm_text: str) -> dict:
    parts = re.split(
        r"(Decline.*?Follow the below steps\.|Refer.*?Follow the below steps\.)",
        ecm_text, flags=re.DOTALL,
    )
    decline_flow: list[dict] = []
    refer_flow: list[dict] = []
    section = None
    for part in parts:
        if "Decline" in part:
            section = "decline"
            continue
        elif "Refer" in part:
            section = "refer"
            continue
        if section == "decline":
            decline_flow = _parse_steps(part)
        elif section == "refer":
            refer_flow = _parse_steps(part)
    return {"DeclineFlow": decline_flow, "ReferFlow": refer_flow}


def parse_address(ecm_text: str) -> dict:
    parts = re.split(
        r"(Refer\s*\(not PASSED\).*?Follow the below steps\.|Refer\s*\+.*?Follow the below steps\.)",
        ecm_text, flags=re.DOTALL,
    )
    decline_flow: list[dict] = []
    refer_flow: list[dict] = []
    section = None
    for part in parts:
        if "not PASSED" in part:
            section = "decline"
            continue
        elif "Refer" in part and section == "decline":
            section = "refer"
            continue
        if section == "decline":
            decline_flow = _parse_steps(part)
        elif section == "refer":
            refer_flow = _parse_steps(part)
    return {"DeclineFlow": decline_flow, "ReferFlow": refer_flow}


def parse_nns(ecm_text: str) -> dict:
    flow = _parse_steps(ecm_text)

    # Country routing from embedded sub-table
    routing: dict[str, str] = {}
    table_match = re.search(r"following table\r+(.*)", ecm_text, re.DOTALL)
    if table_match:
        table_text = table_match.group(1)
        rows = [r.strip() for r in re.split(r"\r\r+", table_text) if r.strip()]
        for row in rows:
            parts = [p.strip() for p in row.split("\r") if p.strip()]
            if len(parts) == 2:
                country, team_str = parts
                if country.lower() in (
                    "country issuing the identity document",
                    "country team",
                ):
                    continue
                team_region = (
                    "APAC" if "APAC" in team_str
                    else ("EMEA" if "EMEA" in team_str else "AMERICAS")
                )
                routing[country] = team_region

    return {"DeclineFlow": flow, "ReferFlow": flow, "CountryRouting": routing}


# ---------------------------------------------------------------------------
# 3. Parse embedded PowerPoint files → document lists
# ---------------------------------------------------------------------------

def parse_pptx_table(pptx_path: str) -> list[dict]:
    """Read THE table from a single-slide PPTX and return list of row dicts."""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    rows: list[dict] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                for row_idx in range(1, len(table.rows)):
                    row_cells = [cell.text.strip() for cell in table.rows[row_idx].cells]
                    rows.append(dict(zip(headers, row_cells)))
    return rows


# ---------------------------------------------------------------------------
# 4. Parse embedded Excel file → NNS Proceed Checklist
# ---------------------------------------------------------------------------

def parse_nns_checklist(xlsx_path: str) -> list[dict]:
    """Read the NNS Proceed Checklist decision matrix from the Excel file."""
    from openpyxl import load_workbook

    wb = load_workbook(xlsx_path)
    ws = wb.active

    # Parse the structured checklist
    # Columns: A=India?, B=Identity Got Exception?, C=Identity Remediation?,
    #          D=Address Got Exception?, E=Address Remediation?, F=Remarks
    rules: list[dict] = []
    for row in ws.iter_rows(min_row=5, max_col=6, values_only=True):
        cells = [str(c).strip() if c else "" for c in row]
        # Skip empty rows
        if not any(cells):
            continue
        # Only parse rows that have a Remarks value
        remark = cells[5]
        if remark and remark not in ("Remarks", ""):
            entry = {
                "india_application": cells[0] if cells[0] else None,
                "identity_exception": cells[1] if cells[1] else None,
                "identity_remediation": cells[2] if cells[2] else None,
                "address_exception": cells[3] if cells[3] else None,
                "address_remediation": cells[4] if cells[4] else None,
                "decision": remark,
            }
            rules.append(entry)
    return rules


# ---------------------------------------------------------------------------
# 5. Main — orchestrate everything
# ---------------------------------------------------------------------------

def main():
    import win32com.client

    doc_path = os.path.abspath(
        r"C:\Users\sharmadeep\kyc-automation\data\SOP - International Account Opening.docx"
    )
    embed_dir = os.path.abspath(
        r"C:\Users\sharmadeep\kyc-automation\data\embedded_objects"
    )
    output_path = os.path.abspath(
        r"C:\Users\sharmadeep\kyc-automation\data\sop_rules\sop_rules_parsed_v2.json"
    )

    print(f"=== SOP Parser v2 ===")
    print(f"Input:  {doc_path}")
    print(f"Embeds: {embed_dir}")
    print(f"Output: {output_path}\n")

    # --- Phase 1: Open Word and extract table + OLE objects ---
    print("Phase 1: Opening SOP document via Word COM...")
    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(doc_path)

    print("  Extracting main SOP table...")
    table_data = extract_sop_table(doc)

    print("  Extracting embedded OLE objects...")
    extracted = extract_embedded_objects(doc, embed_dir)

    doc.Close(False)
    word.Quit()
    print(f"  Word closed. Extracted {len(extracted)} objects.\n")

    # --- Phase 2: Parse SOP table into exception flows ---
    print("Phase 2: Parsing SOP exception flows...")
    rules: dict = {}

    for row_num, row in table_data.items():
        exception_type = row.get(2, "").strip()
        ecm_text = row.get(4, "")

        if "Identity" in exception_type:
            print(f"  Identity exception (row {row_num})")
            rules["Identity"] = parse_identity(ecm_text)
        elif "Address" in exception_type:
            print(f"  Address Proof exception (row {row_num})")
            rules["AddressProof"] = parse_address(ecm_text)
        elif "NNS" in exception_type:
            print(f"  NNS exception (row {row_num})")
            nns = parse_nns(ecm_text)
            rules["NNS"] = {"DeclineFlow": nns["DeclineFlow"], "ReferFlow": nns["ReferFlow"]}
            rules["NNS_CountryRouting"] = nns["CountryRouting"]

    # --- Phase 3: Parse embedded PowerPoint files ---
    print("\nPhase 3: Parsing embedded PowerPoint documents...")

    pptx_mappings = {
        "Acceptable Identity Documents.pptx": ("AcceptableDocuments", "Identity"),
        "Supplement Identity Documents.pptx": ("SupplementGuidelines", "Identity"),
        "Acceptable Address Proof Documents.pptx": ("AcceptableDocuments", "AddressProof"),
        "Supplement Address Proof Documents.pptx": ("SupplementGuidelines", "AddressProof"),
    }

    for label, (key, section) in pptx_mappings.items():
        path = extracted.get(label)
        if path and os.path.isfile(path):
            print(f"  Parsing {label}...")
            rows = parse_pptx_table(path)
            if section in rules:
                rules[section][key] = rows
            else:
                rules[section] = {key: rows}
            print(f"    -> {len(rows)} row(s)")
        else:
            print(f"  MISSING: {label}")

    # --- Phase 4: Parse embedded Excel file ---
    print("\nPhase 4: Parsing NNS Proceed Checklist...")

    checklist_label = "NNS Proceed Checklist.xlsx"
    checklist_path = extracted.get(checklist_label)
    if checklist_path and os.path.isfile(checklist_path):
        checklist_rules = parse_nns_checklist(checklist_path)
        rules["NNS_ProceedChecklist"] = checklist_rules
        print(f"  -> {len(checklist_rules)} rule(s)")
    else:
        print(f"  MISSING: {checklist_label}")

    # --- Phase 5: Note unextracted PDF ---
    pdf_label = "NNS Process.pdf"
    if extracted.get(pdf_label) == "__PDF_NOT_EXTRACTED__":
        rules["NNS_ProcessDocument"] = {
            "status": "not_extracted",
            "reason": "PDF embedded as Acrobat OLE object; requires Acrobat COM or Azure Document Intelligence for extraction",
            "label": pdf_label,
        }
        print(f"\nNote: {pdf_label} could not be extracted (no Acrobat COM available)")

    # --- Write output ---
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(rules, f, indent=2, ensure_ascii=False)

    print(f"\n=== Done! ===")
    print(f"Output: {output_path}")

    # Summary
    print(f"\nSummary:")
    for key in rules:
        if isinstance(rules[key], dict):
            sub_keys = list(rules[key].keys())
            print(f"  {key}: {sub_keys}")
        elif isinstance(rules[key], list):
            print(f"  {key}: {len(rules[key])} entries")
        else:
            print(f"  {key}: {type(rules[key]).__name__}")


if __name__ == "__main__":
    main()
